"""Generate a professional ContractOS Capstone PowerPoint presentation.

Usage:
    python presentation/generate_pptx.py           # writes to presentation/ContractOS_Capstone.pptx
    python presentation/generate_pptx.py out.pptx  # writes to custom path
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Brand palette ──────────────────────────────────────────────────
BG_DARK = RGBColor(0x0A, 0x0A, 0x1A)
SURFACE = RGBColor(0x1E, 0x29, 0x3B)
ACCENT = RGBColor(0x6E, 0xA8, 0xFE)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
RED = RGBColor(0xF8, 0x71, 0x71)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xF0)
DIM_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY = RGBColor(0x47, 0x55, 0x69)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_slide_bg(slide, color: RGBColor = BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, *, font_size=18,
              color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
              font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_title(slide, text, *, top=Inches(0.4), font_size=40):
    _add_text(slide, Inches(0.8), top, Inches(11.7), Inches(0.9), text,
              font_size=font_size, color=WHITE, bold=True,
              alignment=PP_ALIGN.LEFT)


def _add_subtitle(slide, text, *, top=Inches(1.2), font_size=20):
    _add_text(slide, Inches(0.8), top, Inches(11.7), Inches(0.6), text,
              font_size=font_size, color=ACCENT, bold=False,
              alignment=PP_ALIGN.LEFT)


def _add_card(slide, left, top, width, height, title, body, *,
              title_color=ACCENT, border_color=None):
    from pptx.oxml.ns import qn

    shape = slide.shapes.add_shape(
        1, left, top, width, height  # MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x14, 0x17, 0x22)
    shape.line.color.rgb = border_color or RGBColor(0x2D, 0x31, 0x48)
    shape.line.width = Pt(1.5)

    # Round corners via XML
    sp = shape._element.spPr
    prstGeom = sp.find(qn("a:prstGeom"))
    if prstGeom is not None:
        prstGeom.set("prst", "roundRect")

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)

    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(14)
    p_title.font.bold = True
    p_title.font.color.rgb = title_color
    p_title.font.name = "Calibri"

    p_body = tf.add_paragraph()
    p_body.text = body
    p_body.font.size = Pt(11)
    p_body.font.color.rgb = DIM_GRAY
    p_body.font.name = "Calibri"
    p_body.space_before = Pt(6)

    return shape


def _add_table(slide, left, top, width, rows_data, *, col_widths=None):
    rows = len(rows_data)
    cols = len(rows_data[0]) if rows else 0
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.4 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = ACCENT
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE if r == 0 else BG_DARK

    return table_shape


def _add_metric(slide, left, top, number, label, *, width=Inches(2.5)):
    shape = slide.shapes.add_shape(1, left, top, width, Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x14, 0x17, 0x22)
    shape.line.color.rgb = RGBColor(0x2D, 0x31, 0x48)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(12)

    p_num = tf.paragraphs[0]
    p_num.text = str(number)
    p_num.font.size = Pt(36)
    p_num.font.bold = True
    p_num.font.color.rgb = ACCENT
    p_num.font.name = "Calibri"
    p_num.alignment = PP_ALIGN.CENTER

    p_label = tf.add_paragraph()
    p_label.text = label
    p_label.font.size = Pt(11)
    p_label.font.color.rgb = DIM_GRAY
    p_label.font.name = "Calibri"
    p_label.alignment = PP_ALIGN.CENTER


def _add_traffic_dot(slide, left, top, label, sublabel, color):
    circle = slide.shapes.add_shape(9, left, top, Inches(0.9), Inches(0.9))  # OVAL
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    _add_text(slide, left - Inches(0.3), top + Inches(1.05), Inches(1.5), Inches(0.6),
              f"{label}\n{sublabel}", font_size=10, color=DIM_GRAY,
              alignment=PP_ALIGN.CENTER)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # Blank

    # ── SLIDE 1: Title ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_text(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.5),
              "ContractOS", font_size=60, color=WHITE, bold=True,
              alignment=PP_ALIGN.CENTER)
    _add_text(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(0.7),
              "The Operating System for Contract Intelligence",
              font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)
    _add_text(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.5),
              "Capstone Project Presentation",
              font_size=16, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.4),
              "Python  •  FastAPI  •  FAISS  •  Anthropic Claude  •  MCP  •  D3.js",
              font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 2: The Problem ────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "The Problem")
    _add_text(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5),
              "Contracts encode obligations, risks, and rights implicitly.",
              font_size=18, color=LIGHT_GRAY)

    cards = [
        ("Buried in Tables", "SLA penalties, pricing tiers, and delivery schedules hidden in schedule annexes that nobody reads until it's too late."),
        ("Scattered Clauses", '"Subject to Section 5.2" — cross-references create a web of dependencies that humans struggle to trace.'),
        ("Defined Terms", '"Service Provider" means TechServe... but also their subcontractors? Definitions transform meaning silently.'),
        ("No Provenance", 'When an AI says "the liability cap is $14M" — where exactly does it say that? Can you trust it?'),
    ]
    for i, (title, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        _add_card(slide,
                  Inches(0.8 + col * 6.2), Inches(2.1 + row * 2.3),
                  Inches(5.8), Inches(2.0),
                  title, body)

    # ── SLIDE 3: Market Gap ─────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "What Existing Tools Get Wrong")
    _add_table(slide, Inches(1.5), Inches(1.6), Inches(10.3), [
        ["Approach", "Problem"],
        ["Keyword Search", 'Misses "shall indemnify" when searching for "liability"'],
        ["Document Chatbots", "Hallucinate answers with no source grounding"],
        ["Flat Embedding Search", "Treats a 50-page contract as a text blob"],
        ["Manual Review", "8+ hours per contract, inconsistent, doesn't scale"],
    ], col_widths=[Inches(3), Inches(7.3)])
    _add_text(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.8),
              "They all fail at the same thing: separating what a contract says from what it means from what someone thinks about it.",
              font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 4: Truth Model ────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "Our Solution: The Truth Model")
    _add_text(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5),
              "A strict four-layer epistemological model for contract intelligence.",
              font_size=18, color=LIGHT_GRAY)

    truth_cards = [
        ("FACT", "Directly grounded in contract text. Immutable.", '"Term: 24 months"', GREEN),
        ("BINDING", "Explicit semantic mapping. Scoped.", '"Company" = Acme Corp', ACCENT),
        ("INFERENCE", "Derived claim with confidence. Revisable.", "Expires Dec 2025 (0.92)", YELLOW),
        ("OPINION", "Contextual judgment. Never persisted as truth.", "Missing force majeure = risk", RED),
    ]
    for i, (title, desc, example, color) in enumerate(truth_cards):
        x = Inches(0.6 + i * 3.15)
        _add_card(slide, x, Inches(2.2), Inches(2.9), Inches(3.0),
                  title, f"{desc}\n\n{example}", title_color=color, border_color=color)

    _add_text(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.5),
              "Every answer is auditable. Every claim is traceable to source text.",
              font_size=14, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 5: Architecture ───────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "Architecture")

    layers = [
        ("INTERACTION LAYER", "Browser Copilot  •  REST API (34 endpoints + SSE)  •  TrustGraph Viz  •  MCP Server  •  Cross-Contract  •  Reports  •  Streaming", ACCENT),
        ("AGENT LAYER", "DocumentAgent (Semantic Q&A + FAISS + LLM)  •  ComplianceAgent (Playbook G/Y/R)  •  NDATriageAgent (10-point)  •  DraftAgent (Redline)", PURPLE),
        ("EXTRACTION TOOLS", "FactExtractor (0 LLM)  •  BindingResolver  •  ClauseClassifier (18 types)  •  AliasDetector  •  CrossRefExtractor  •  FactDiscovery (LLM)  •  DocParser", GREEN),
        ("DATA FABRIC", "TrustGraph (SQLite WAL)  •  EmbeddingIndex (FAISS MiniLM-L6-v2 384-dim)  •  WorkspaceStore  •  LLM Provider (Anthropic Claude)", YELLOW),
        ("TRUTH MODEL", "FACT (immutable, char offsets)  •  BINDING (semantic mapping)  •  INFERENCE (confidence scored)  •  OPINION (never persisted)", RED),
    ]
    for i, (name, desc, color) in enumerate(layers):
        y = Inches(1.3 + i * 1.15)
        shape = slide.shapes.add_shape(1, Inches(0.8), y, Inches(11.7), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(14)
        tf.margin_top = Pt(8)

        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = "Calibri"

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = DIM_GRAY
        p2.font.name = "Calibri"
        p2.space_before = Pt(4)

    # Provenance arrow label
    _add_text(slide, Inches(11.5), Inches(3.0), Inches(1.5), Inches(1.5),
              "PROVENANCE\nCHAIN  ↑", font_size=10, color=RED,
              bold=True, alignment=PP_ALIGN.CENTER)

    _add_text(slide, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3),
              "Downward-only data flow  |  Every output is typed  |  Agents are stateless  |  Interaction layer never reasons",
              font_size=9, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 6: Three-Phase Pipeline ───────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "Three-Phase Extraction Pipeline")

    phases = [
        ("1", "Pattern Extraction", "Regex + structural parsing\nDefinitions, dates, amounts\nZero LLM calls", GREEN),
        ("2", "Semantic Indexing", "FAISS + MiniLM-L6-v2\n384-dim vectors\nPer-document indexes", ACCENT),
        ("3", "LLM Discovery", "Anthropic Claude\nHidden obligations, risks\nCross-clause implications", PURPLE),
    ]
    for i, (num, title, desc, color) in enumerate(phases):
        x = Inches(1.0 + i * 4.0)
        circle = slide.shapes.add_shape(9, x, Inches(2.0), Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BG_DARK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        _add_text(slide, x - Inches(0.5), Inches(3.0), Inches(1.8), Inches(0.4),
                  title, font_size=16, color=WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER)
        _add_text(slide, x - Inches(0.8), Inches(3.5), Inches(2.4), Inches(1.5),
                  desc, font_size=12, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

        if i < 2:
            _add_text(slide, x + Inches(1.6), Inches(2.1), Inches(1.5), Inches(0.7),
                      "→", font_size=36, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    _add_text(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.6),
              "Phases 1 & 2 are deterministic and reproducible. Phase 3 is AI-augmented — it finds what patterns miss.",
              font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 7: Upload & Semantic Q&A ──────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "Upload & Semantic Q&A")
    _add_subtitle(slide, "Upload any PDF or DOCX. Ask questions. Get grounded answers.")

    steps = [
        'User: "Does this contract indemnify the buyer for data breach?"',
        "",
        "ContractOS:",
        "  1. Extracts 520+ facts from the document",
        "  2. Resolves 27 bindings (\"Service Provider\" → \"TechServe Solutions\")",
        "  3. Searches FAISS index for indemnity-related facts",
        "  4. Reasons with Claude, grounded in extracted evidence",
        "  5. Returns answer with confidence: very_high, provenance chain, clickable source highlighting",
    ]
    shape = slide.shapes.add_shape(1, Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x10, 0x14, 0x25)
    shape.line.color.rgb = RGBColor(0x2D, 0x31, 0x48)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.margin_top = Pt(16)
    for i, line in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = "Calibri"
        p.font.color.rgb = ACCENT if line.startswith("User:") or line.startswith("ContractOS:") else LIGHT_GRAY
        p.font.bold = line.startswith("User:") or line.startswith("ContractOS:")
        p.space_before = Pt(4)

    # ── SLIDE 8: Playbook Compliance ────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "Playbook Compliance Review")
    _add_subtitle(slide, "Compare contracts against your organization's standard positions.")

    _add_traffic_dot(slide, Inches(3.5), Inches(2.2), "GREEN", "Meets standard", GREEN)
    _add_traffic_dot(slide, Inches(6.2), Inches(2.2), "YELLOW", "Deviates — negotiate", YELLOW)
    _add_traffic_dot(slide, Inches(8.9), Inches(2.2), "RED", "Missing or unacceptable", RED)

    review_cards = [
        ("10 Clause Types", "Liability, indemnification, confidentiality, termination, IP, data protection, force majeure, governing law, assignment, payment"),
        ("Automated Redlines", "LLM-generated alternative language with rationale and fallback positions for every YELLOW/RED finding"),
        ("Risk Matrix", "5×5 Severity × Likelihood scoring with aggregate risk profile and tier distribution"),
        ("Real-Time Streaming", "SSE delivers each clause evaluation as it happens — watch the AI reason in real time"),
    ]
    for i, (title, body) in enumerate(review_cards):
        col = i % 2
        row = i // 2
        _add_card(slide, Inches(0.8 + col * 6.2), Inches(4.2 + row * 1.7),
                  Inches(5.8), Inches(1.5), title, body)

    # ── SLIDE 9: NDA Triage ─────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "NDA Triage Screening")
    _add_subtitle(slide, "Automated 10-point checklist with hybrid pattern + LLM evaluation.")

    _add_table(slide, Inches(2.0), Inches(2.0), Inches(9.3), [
        ["#", "Check", "Method"],
        ["1", "Mutual obligations", "Pattern + LLM"],
        ["2", "Standard carveouts", "Pattern matching"],
        ["3", "Term duration", "Pattern extraction"],
        ["4", "Governing law", "Pattern extraction"],
        ["5-6", "IP assignment • Non-compete", "LLM analysis"],
        ["7-8", "Indemnification • Problematic provisions", "Pattern + LLM"],
        ["9-10", "Data protection • Agreement structure", "Structural + LLM"],
    ], col_widths=[Inches(1), Inches(5), Inches(3.3)])

    _add_text(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.5),
              "GREEN = auto-approve  •  YELLOW = expedited review  •  RED = full legal review",
              font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 10: Risk & Obligations ────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "Risk Memo & Obligation Extraction")

    risk_cards = [
        ("Risk Assessment Memo", "Executive summary, key risks with severity/likelihood scoring, missing protections, prioritized recommendations, escalation items."),
        ("Obligation Extraction", "All contractual obligations categorized: Affirmative (must do), Negative (must not do), Conditional (if X then Y) — with deadlines."),
        ("Hidden Fact Discovery", "LLM-powered analysis surfacing implicit obligations, liability exposure, unstated assumptions, and ambiguous terms."),
        ("Clause Gap Analysis", "Identifies mandatory facts missing from each clause type — what should be there but isn't."),
    ]
    for i, (title, body) in enumerate(risk_cards):
        col = i % 2
        row = i // 2
        _add_card(slide, Inches(0.8 + col * 6.2), Inches(1.6 + row * 2.6),
                  Inches(5.8), Inches(2.3), title, body)

    # ── SLIDE 11: MCP Integration ───────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "MCP Server Integration")
    _add_subtitle(slide, "Full contract intelligence as an MCP server — works in Cursor, Claude Desktop, Claude Code.")

    tools = [
        "upload_contract", "ask_question", "review_against_playbook",
        "triage_nda", "discover_hidden_facts", "extract_obligations",
        "generate_risk_memo", "compare_clauses", "search_contracts",
        "get_clause_gaps", "generate_report", "load_sample_contract",
    ]
    for i, tool in enumerate(tools):
        col = i % 4
        row = i // 4
        _add_card(slide, Inches(0.6 + col * 3.15), Inches(2.2 + row * 1.5),
                  Inches(2.9), Inches(1.2), tool, "", title_color=PURPLE,
                  border_color=RGBColor(0x3D, 0x2E, 0x5E))

    _add_text(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4),
              "+ 10 Resources (read-only data)  •  5 Prompts (reusable workflows)  •  stdio + HTTP transport",
              font_size=11, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 12: Demo Scenario 1 ──────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "Demo Scenario 1")
    _add_subtitle(slide, "Complex Procurement Framework (GBP 85M)")

    scenario1 = (
        'The Situation: "Due to a recent interstate transport strike, a lorry carrying raw materials '
        "from our supplier GlobalSource got delayed. The goods arrived at our godown 10 days late. "
        'We had to procure from a local vendor at higher cost."\n\n'
        'The Question: "What are my legal options? Can I reject the late delivery, cancel the order, '
        'and claim my money back? Am I entitled to liquidated damages?"'
    )
    shape = slide.shapes.add_shape(1, Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x10, 0x14, 0x25)
    shape.line.color.rgb = RGBColor(0x2D, 0x31, 0x48)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.margin_top = Pt(16)
    tf.paragraphs[0].text = scenario1
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
    tf.paragraphs[0].font.name = "Calibri"

    _add_text(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.6),
              "ContractOS analyses force majeure clauses, liquidated damages schedules, delivery obligations, "
              "and rejection rights — all with provenance back to specific contract sections.",
              font_size=13, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 13: Demo Scenario 2 ──────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "Demo Scenario 2")
    _add_subtitle(slide, "IT Outsourcing Agreement ($47.5M)")

    scenario2 = (
        'The Situation: "Our outsourcing vendor TechServe experienced a data breach affecting 50,000 '
        "customer records including PII. They took 3 days to notify us — violating the 24-hour "
        'breach notification requirement."\n\n'
        'The Question: "Can we terminate for cause? What damages can we claim? '
        'Does the liability cap apply to data breaches?"'
    )
    shape = slide.shapes.add_shape(1, Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x10, 0x14, 0x25)
    shape.line.color.rgb = RGBColor(0x2D, 0x31, 0x48)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.margin_top = Pt(16)
    tf.paragraphs[0].text = scenario2
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
    tf.paragraphs[0].font.name = "Calibri"

    _add_text(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.6),
              "ContractOS traces through SLA tables, breach notification clauses, liability carve-outs, "
              "and indemnification provisions — showing exactly which sections apply.",
              font_size=13, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 14: Cross-Contract Intelligence ───────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "Cross-Contract Intelligence")
    _add_subtitle(slide, "Compare provisions across multiple contracts simultaneously.")

    _add_table(slide, Inches(0.8), Inches(2.0), Inches(11.7), [
        ["Clause", "Procurement (GBP 85M)", "IT Outsourcing ($47.5M)"],
        ["Force Majeure", "90-day threshold, mutual", "60-day threshold, pandemic-specific"],
        ["Liability Cap", "150% of framework value", "150% annual fees ($14.25M)"],
        ["Termination Notice", "90 days (buyer only)", "180 days (mutual)"],
        ["Data Breach", "Not addressed", "24-hour notification, unlimited liability"],
        ["Governing Law", "England & Wales", "New York"],
    ], col_widths=[Inches(2.5), Inches(4.6), Inches(4.6)])

    _add_text(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.8),
              'Ask: "Which contract has stronger force majeure protection?" — and get a grounded, '
              "comparative answer with provenance from both documents.",
              font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 15: By the Numbers ────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "By the Numbers")

    metrics = [
        ("794", "Tests Passing"), ("34", "API Endpoints"),
        ("13", "MCP Tools"), ("50", "Real NDAs Tested"),
        ("9.8K", "Lines of Production Code"), ("14.8K", "Lines of Test Code"),
        ("63", "Python Modules"), ("13", "Dev Phases Complete"),
    ]
    for i, (num, label) in enumerate(metrics):
        col = i % 4
        row = i // 4
        _add_metric(slide, Inches(0.8 + col * 3.1), Inches(1.8 + row * 2.2), num, label)

    # ── SLIDE 16: Tech Stack ────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "Technology Stack")

    tech = [
        ("Python 3.12+", "Backend & CLI"), ("FastAPI", "Async REST API"),
        ("Anthropic Claude", "LLM Reasoning"), ("FAISS", "Vector Search"),
        ("sentence-transformers", "MiniLM-L6-v2 Embeddings"), ("SQLite (WAL)", "TrustGraph Storage"),
        ("python-docx", "DOCX Parsing"), ("PyMuPDF", "PDF Parsing"),
        ("SSE", "Real-Time Streaming"), ("MCP (FastMCP)", "Model Context Protocol"),
        ("D3.js", "Graph Visualization"), ("Docker", "Single-Container Deploy"),
    ]
    for i, (name, desc) in enumerate(tech):
        col = i % 4
        row = i // 4
        _add_card(slide, Inches(0.6 + col * 3.15), Inches(1.6 + row * 1.8),
                  Inches(2.9), Inches(1.5), name, desc)

    # ── SLIDE 17: TDD & Quality ─────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "Test-Driven Development")
    _add_subtitle(slide, "The entire system was built TDD — tests written before code.")

    test_bars = [
        ("Unit", 539, GREEN, Inches(7.5)),
        ("Integration", 156, ACCENT, Inches(4.5)),
        ("Benchmark", 61, PURPLE, Inches(2.5)),
        ("Contract", 27, YELLOW, Inches(1.5)),
    ]
    for i, (label, count, color, bar_w) in enumerate(test_bars):
        y = Inches(2.4 + i * 0.9)
        _add_text(slide, Inches(1.5), y, Inches(2.0), Inches(0.5),
                  label, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)
        bar = slide.shapes.add_shape(1, Inches(3.8), y + Inches(0.05), bar_w, Inches(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        tf = bar.text_frame
        tf.paragraphs[0].text = f"  {count} tests"
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BG_DARK
        tf.paragraphs[0].font.name = "Calibri"
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    _add_card(slide, Inches(1.5), Inches(5.8), Inches(4.5), Inches(1.2),
              "1.5x Test-to-Code Ratio",
              "14,800 lines of tests vs 9,800 lines of production code")
    _add_card(slide, Inches(7.0), Inches(5.8), Inches(4.5), Inches(1.2),
              "50 Real NDAs",
              "ContractNLI dataset from Stanford NLP — corporate, M&A, government, SEC filings")

    # ── SLIDE 18: Live Demo ─────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_text(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(1.0),
              "Live Demo", font_size=48, color=WHITE, bold=True,
              alignment=PP_ALIGN.CENTER)
    _add_text(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(0.6),
              "Watch ContractOS in Action", font_size=22, color=ACCENT,
              alignment=PP_ALIGN.CENTER)

    acts = [
        ("Act 1", "Procurement Framework\nLorry strike scenario\nForce majeure & damages", GREEN),
        ("Act 2", "IT Outsourcing\nData breach scenario\nMulti-doc workspace", ACCENT),
        ("Act 3", "Cross-Contract Intel\nCompare & Q&A\nMerged TrustGraph", PURPLE),
    ]
    for i, (title, desc, color) in enumerate(acts):
        _add_card(slide, Inches(1.5 + i * 3.8), Inches(3.2), Inches(3.3), Inches(2.5),
                  title, desc, title_color=color, border_color=color)

    # ── SLIDE 19: What's Next ───────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_title(slide, "What's Next")

    roadmap = [
        ("✓", "Phase 1-13: Complete", "Full extraction, Q&A, playbook review, NDA triage, MCP server, streaming, cross-contract intelligence", GREEN),
        ("→", "Contract Families", "Master agreements + amendments as DAGs with precedence rules", ACCENT),
        ("→", "Multi-Tenant Workspaces", "Organization-scoped playbooks, role-based access, team collaboration", ACCENT),
        ("→", "Clause Library", "Pre-approved clause templates with version history and approval workflows", ACCENT),
        ("→", "Portfolio Analytics", "Risk heatmaps across entire contract portfolios, renewal tracking, obligation calendars", ACCENT),
        ("→", "Word/Outlook Add-in", "In-document copilot for Microsoft 365 with real-time review", ACCENT),
    ]
    for i, (icon, title, desc, color) in enumerate(roadmap):
        y = Inches(1.5 + i * 0.9)
        dot = slide.shapes.add_shape(9, Inches(1.0), y + Inches(0.08), Inches(0.25), Inches(0.25))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        _add_text(slide, Inches(1.5), y, Inches(3.0), Inches(0.4),
                  title, font_size=14, color=WHITE, bold=True)
        _add_text(slide, Inches(1.5), y + Inches(0.35), Inches(10.0), Inches(0.4),
                  desc, font_size=11, color=DIM_GRAY)

    # ── SLIDE 20: Thank You ─────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_text(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
              "Thank You", font_size=60, color=WHITE, bold=True,
              alignment=PP_ALIGN.CENTER)
    _add_text(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
              "Don't read contracts. Understand them.",
              font_size=28, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.6),
              "ContractOS — The Operating System for Contract Intelligence",
              font_size=16, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(0.4),
              "794 tests  •  34 API endpoints  •  13 MCP tools  •  50 real NDAs tested",
              font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.4),
              "Built with Python  •  FastAPI  •  FAISS  •  Anthropic Claude  •  SQLite  •  Docker",
              font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    return prs


def main():
    out_path = sys.argv[1] if len(sys.argv) > 1 else str(
        Path(__file__).parent / "ContractOS_Capstone.pptx"
    )
    prs = build_presentation()
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
